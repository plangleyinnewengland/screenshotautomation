"""
Camtasia Project Generator - Create Camtasia projects from image folders
Supports Camtasia 2025 (version 25.x) JSON format (.tscproj)
"""

import os
import sys
import json
import uuid
import csv
import wave
from pathlib import Path
from datetime import datetime
from typing import List, Optional, Dict
import time

try:
    from PIL import Image
except ImportError:
    print("PIL not available - image dimensions will use defaults")
    Image = None


def generate_guid():
    """Generate a GUID in Camtasia format."""
    return str(uuid.uuid4()).upper()


def get_image_info(image_path: Path) -> Dict:
    """Get image dimensions and format."""
    if Image:
        try:
            with Image.open(image_path) as img:
                return {
                    'width': img.width,
                    'height': img.height,
                    'format': img.format or 'PNG'
                }
        except:
            pass
    # Default to 1920x1080 if we can't read
    return {'width': 1920, 'height': 1080, 'format': 'PNG'}


def get_images_from_directory(directory: str) -> List[Path]:
    """Get all image files from a directory, sorted by name."""
    directory = Path(directory)
    if not directory.exists():
        raise FileNotFoundError(f"Directory not found: {directory}")
    
    image_extensions = {'.png', '.jpg', '.jpeg', '.bmp', '.gif', '.tiff', '.webp'}
    
    image_files = sorted([
        f for f in directory.iterdir()
        if f.is_file() and f.suffix.lower() in image_extensions
    ])
    
    return image_files


def get_audio_from_directory(directory: str) -> List[Path]:
    """Get all supported audio files from a directory, sorted by name."""
    directory = Path(directory)
    if not directory.exists():
        raise FileNotFoundError(f"Directory not found: {directory}")

    audio_extensions = {'.wav', '.mp3', '.m4a', '.aac', '.wma'}
    audio_files = sorted([
        f for f in directory.iterdir()
        if f.is_file() and f.suffix.lower() in audio_extensions
    ])
    return audio_files


def get_audio_duration_seconds(audio_path: Path, fallback_seconds: float = 5.0) -> float:
    """Best-effort audio duration. Uses wave for WAV, otherwise falls back."""
    try:
        if audio_path.suffix.lower() == '.wav':
            with wave.open(str(audio_path), 'rb') as wav_file:
                frame_count = wav_file.getnframes()
                frame_rate = wav_file.getframerate()
                if frame_rate > 0:
                    return frame_count / float(frame_rate)
    except Exception:
        pass

    return fallback_seconds


def _extract_speaker_notes_text(slide) -> str:
    """Extract plain text speaker notes from a slide."""
    try:
        notes_slide = slide.notes_slide
        text = notes_slide.notes_text_frame.text or ""
        lines = [line.strip() for line in text.splitlines() if line.strip()]
        return "\n".join(lines)
    except Exception:
        return ""


def create_audiate_script_from_pptx(
    pptx_path: str,
    output_path: str = None,
    default_voice: str = "Andrew (US)"
) -> str:
    """
    Create an Audiate-friendly CSV script from PowerPoint speaker notes.

    Columns:
      - Slide
      - Voice
      - SpeakerNotes
      - SuggestedAudioFile
    """
    try:
        from pptx import Presentation
    except ImportError as exc:
        raise RuntimeError("python-pptx is required to extract speaker notes") from exc

    pptx_file = Path(pptx_path)
    if not pptx_file.exists():
        raise FileNotFoundError(f"PowerPoint file not found: {pptx_file}")

    if output_path is None:
        output_path = str(pptx_file.with_name(f"{pptx_file.stem}_audiate_script.csv"))

    presentation = Presentation(str(pptx_file))

    rows = []
    for slide_index, slide in enumerate(presentation.slides, 1):
        notes_text = _extract_speaker_notes_text(slide)
        if not notes_text:
            continue

        rows.append({
            "Slide": slide_index,
            "Voice": default_voice,
            "SpeakerNotes": notes_text,
            "SuggestedAudioFile": f"{slide_index:04d}.wav",
        })

    if not rows:
        raise ValueError("No speaker notes found in the selected PowerPoint file")

    with open(output_path, 'w', newline='', encoding='utf-8') as csv_file:
        writer = csv.DictWriter(
            csv_file,
            fieldnames=["Slide", "Voice", "SpeakerNotes", "SuggestedAudioFile"]
        )
        writer.writeheader()
        writer.writerows(rows)

    print(f"Created Audiate script: {output_path}")
    print(f"  - Notes rows exported: {len(rows)}")
    print(f"  - Default narrator voice: {default_voice}")

    return output_path


def _compact_iso_timestamp(dt: datetime = None) -> str:
    """Generate compact timestamp for Camtasia metadata fields."""
    if dt is None:
        dt = datetime.now()
    return dt.strftime("%Y%m%dT%H%M%S")


CAMTASIA_EDIT_RATE = 705600000


def _time_added_timestamp(dt: datetime = None) -> str:
    """Generate Camtasia sourceBin timestamp with microseconds."""
    if dt is None:
        dt = datetime.now()
    return dt.strftime("%Y%m%dT%H%M%S.%f")


def _seconds_to_ticks(seconds: float) -> int:
    return max(1, int(seconds * CAMTASIA_EDIT_RATE))


def _file_last_mod_timestamp(path: Path) -> str:
    try:
        ts = datetime.fromtimestamp(path.stat().st_mtime)
        return _compact_iso_timestamp(ts)
    except Exception:
        return _compact_iso_timestamp()


def _to_project_src(asset_path: Path, output_project_path: Path) -> str:
    """Prefer relative source paths when possible, otherwise absolute path."""
    try:
        return str(asset_path.resolve().relative_to(output_project_path.parent.resolve())).replace('/', '\\')
    except Exception:
        return str(asset_path.resolve())


def _default_track_attributes() -> List[Dict]:
    return [
        {
            "ident": "",
            "audioMuted": False,
            "videoHidden": False,
            "magnetic": False,
            "matte": 0,
            "solo": False,
            "metadata": {"IsLocked": "False", "WinTrackHeight": "56"},
        },
        {
            "ident": "",
            "audioMuted": False,
            "videoHidden": False,
            "magnetic": False,
            "matte": 0,
            "solo": False,
            "metadata": {"IsLocked": "False", "WinTrackHeight": "56"},
        },
    ]


def _default_caption_attributes() -> Dict:
    return {
        "enabled": True,
        "fontName": "Arial",
        "fontSize": 64,
        "backgroundColor": [0, 0, 0, 191],
        "foregroundColor": [255, 255, 255, 255],
        "lang": "en",
        "alignment": 0,
        "defaultFontSize": True,
        "opacity": 0.5,
        "backgroundEnabled": True,
        "backgroundOnlyAroundText": True,
    }


def _get_audio_track_info(audio_path: Path, fallback_seconds: float = 5.0) -> Dict:
    """Best-effort extraction of audio track technical metadata."""
    info = {
        "duration": fallback_seconds,
        "sample_rate": 44100,
        "bit_depth": 16,
        "num_channels": 2,
        "sample_count": int(44100 * fallback_seconds),
    }
    try:
        if audio_path.suffix.lower() == '.wav':
            with wave.open(str(audio_path), 'rb') as wav_file:
                frames = wav_file.getnframes()
                sample_rate = wav_file.getframerate() or 44100
                channels = wav_file.getnchannels() or 1
                sample_width = wav_file.getsampwidth() or 2
                duration = frames / float(sample_rate) if sample_rate else fallback_seconds
                info.update({
                    "duration": duration,
                    "sample_rate": sample_rate,
                    "bit_depth": sample_width * 8,
                    "num_channels": channels,
                    "sample_count": frames,
                })
        else:
            duration = get_audio_duration_seconds(audio_path, fallback_seconds=fallback_seconds)
            info["duration"] = duration
            info["sample_count"] = int(info["sample_rate"] * duration)
    except Exception:
        pass
    return info


def _build_project_base(
    title: str,
    project_width: int,
    project_height: int,
    framerate: float,
    audio_sample_rate: int,
) -> Dict:
    now = datetime.now()
    return {
        "title": title,
        "description": "",
        "author": "",
        "targetLoudness": -18.0,
        "shouldApplyLoudnessNormalization": True,
        "videoFormatFrameRate": int(round(framerate)),
        "audioFormatSampleRate": int(audio_sample_rate),
        "allowSubFrameEditing": False,
        "width": float(project_width),
        "height": float(project_height),
        "version": "9.0",
        "editRate": CAMTASIA_EDIT_RATE,
        "authoringClientName": {
            "name": "Camtasia",
            "platform": "Windows",
            "version": "25.2.2",
        },
        "sourceBin": [],
        "timeline": {
            "id": 2,
            "sceneTrack": {
                "scenes": [{
                    "csml": {
                        "tracks": [
                            {"trackIndex": 0, "medias": []},  # audio track
                            {"trackIndex": 1, "medias": []},  # image track
                        ]
                    }
                }]
            },
            "parameters": {},
            "trackAttributes": _default_track_attributes(),
            "captionAttributes": _default_caption_attributes(),
            "gain": 1.0,
            "legacyAttenuateAudioMix": False,
            "backgroundColor": [0, 0, 0, 255],
        },
        "metadata": {
            "Date": now.strftime("%Y-%m-%d %I:%M:%S %p"),
            "IsStandalone": True,
            "Language": "ENU",
            "ProfileName": "",
            "Title": title,
            "audioNarrationNotes": "",
            "calloutStyle": "Basic",
        },
    }


def create_camtasia_project_json(
    image_folder: str,
    output_path: str = None,
    image_duration: float = 5.0,
    project_width: int = 1920,
    project_height: int = 1080,
    framerate: float = 30.0
) -> str:
    """Create a Camtasia project from images using Camtasia's native schema."""
    image_folder = Path(image_folder)
    images = get_images_from_directory(image_folder)
    
    if not images:
        raise ValueError(f"No images found in {image_folder}")
    
    if output_path is None:
        output_path = str(image_folder / f"{image_folder.name}_project.tscproj")
    
    output_project = Path(output_path)
    seconds_per_image = max(0.1, float(image_duration))

    project = _build_project_base(
        title=output_project.stem,
        project_width=project_width,
        project_height=project_height,
        framerate=framerate,
        audio_sample_rate=44100,
    )

    scene_tracks = project["timeline"]["sceneTrack"]["scenes"][0]["csml"]["tracks"]
    image_track_medias = scene_tracks[1]["medias"]

    source_id = 2
    media_id = 100
    start_ticks = 0

    for index, img_path in enumerate(images, 1):
        img_info = get_image_info(img_path)
        src = _to_project_src(img_path, output_project)

        project["sourceBin"].append({
            "id": source_id,
            "src": src,
            "rect": [0, 0, int(img_info['width']), int(img_info['height'])],
            "lastMod": _file_last_mod_timestamp(img_path),
            "loudnessNormalization": True,
            "sourceTracks": [{
                "range": [0, 1],
                "type": 1,
                "editRate": 10000000,
                "trackRect": [0, 0, int(img_info['width']), int(img_info['height'])],
                "sampleRate": 0,
                "bitDepth": 24,
                "numChannels": 0,
                "integratedLUFS": 100.0,
                "peakLevel": -1.0,
                "tag": 0,
                "metaData": f"{img_path.name};",
                "parameters": {},
            }],
            "metadata": {"timeAdded": _time_added_timestamp()},
        })

        duration_ticks = _seconds_to_ticks(seconds_per_image)
        image_track_medias.append({
            "id": media_id,
            "_type": "IMFile",
            "src": source_id,
            "trackNumber": 0,
            "trimStartSum": 0,
            "attributes": {"ident": f"slide_{index}"},
            "parameters": {
                "geometryCrop0": 0.0,
                "geometryCrop1": 0.0,
                "geometryCrop2": 0.0,
                "geometryCrop3": 0.0,
            },
            "effects": [],
            "start": start_ticks,
            "duration": duration_ticks,
            "mediaStart": 0,
            "mediaDuration": 1,
            "scalar": 1,
            "metadata": {
                "audiateLinkedSession": "",
                "clipSpeedAttribute": {"type": "bool", "value": False},
                "colorAttribute": {"type": "color", "value": [0, 0, 0, 0]},
                "default-height": {"type": "double", "value": float(project_height)},
                "default-translation0": {"type": "double", "value": 0.0},
                "default-translation1": {"type": "double", "value": 0.0},
                "default-translation2": {"type": "double", "value": 0.0},
                "default-width": {"type": "double", "value": float(project_width)},
            },
            "animationTracks": {},
        })

        start_ticks += duration_ticks
        source_id += 1
        media_id += 1
    
    # Write to file
    with open(output_path, 'w', encoding='utf-8-sig') as f:
        json.dump(project, f, indent=2)
    
    print(f"Created Camtasia project: {output_path}")
    print(f"  - {len(images)} images imported")
    print(f"  - Duration per image: {image_duration}s")
    print(f"  - Total duration: {len(images) * image_duration}s")
    print(f"  - Resolution: {project_width}x{project_height}")
    print(f"  - Framerate: {int(round(framerate))} fps")
    
    return output_path


def create_camtasia_project_with_audio_json(
    image_folder: str,
    audio_folder: str,
    output_path: str = None,
    min_image_duration: float = 2.0,
    project_width: int = 1920,
    project_height: int = 1080,
    framerate: float = 30.0
) -> str:
    """Create a Camtasia project where each image is aligned to an audio file."""
    image_folder = Path(image_folder)
    audio_folder = Path(audio_folder)

    images = get_images_from_directory(image_folder)
    audio_files = get_audio_from_directory(audio_folder)

    if not images:
        raise ValueError(f"No images found in {image_folder}")
    if not audio_files:
        raise ValueError(f"No audio files found in {audio_folder}")

    if output_path is None:
        output_path = str(image_folder / f"{image_folder.name}_audiate_project.tscproj")

    output_project = Path(output_path)
    pair_count = min(len(images), len(audio_files))
    if pair_count == 0:
        raise ValueError("Need at least one matching image/audio pair")

    audio_infos = [_get_audio_track_info(a, fallback_seconds=min_image_duration) for a in audio_files[:pair_count]]
    default_sample_rate = audio_infos[0]["sample_rate"] if audio_infos else 44100
    project = _build_project_base(
        title=output_project.stem,
        project_width=project_width,
        project_height=project_height,
        framerate=framerate,
        audio_sample_rate=default_sample_rate,
    )

    scene_tracks = project["timeline"]["sceneTrack"]["scenes"][0]["csml"]["tracks"]
    audio_track_medias = scene_tracks[0]["medias"]
    image_track_medias = scene_tracks[1]["medias"]

    source_id = 2
    media_id = 100
    start_ticks = 0

    image_source_ids = []
    audio_source_ids = []

    # Add image and audio sources
    for i in range(pair_count):
        img_path = images[i]
        audio_path = audio_files[i]
        img_info = get_image_info(img_path)
        audio_info = audio_infos[i]

        img_src = _to_project_src(img_path, output_project)
        aud_src = _to_project_src(audio_path, output_project)

        image_source_id = source_id
        project["sourceBin"].append({
            "id": image_source_id,
            "src": img_src,
            "rect": [0, 0, int(img_info['width']), int(img_info['height'])],
            "lastMod": _file_last_mod_timestamp(img_path),
            "loudnessNormalization": True,
            "sourceTracks": [{
                "range": [0, 1],
                "type": 1,
                "editRate": 10000000,
                "trackRect": [0, 0, int(img_info['width']), int(img_info['height'])],
                "sampleRate": 0,
                "bitDepth": 24,
                "numChannels": 0,
                "integratedLUFS": 100.0,
                "peakLevel": -1.0,
                "tag": 0,
                "metaData": f"{img_path.name};",
                "parameters": {},
            }],
            "metadata": {"timeAdded": _time_added_timestamp()},
        })
        source_id += 1

        audio_source_id = source_id
        project["sourceBin"].append({
            "id": audio_source_id,
            "src": aud_src,
            "rect": [0, 0, 0, 0],
            "lastMod": _file_last_mod_timestamp(audio_path),
            "loudnessNormalization": True,
            "sourceTracks": [{
                "range": [0, int(audio_info["sample_count"])],
                "type": 2,
                "editRate": int(audio_info["sample_rate"]),
                "trackRect": [0, 0, 0, 0],
                "sampleRate": int(audio_info["sample_rate"]),
                "bitDepth": int(audio_info["bit_depth"]),
                "numChannels": int(audio_info["num_channels"]),
                "integratedLUFS": -18.0,
                "peakLevel": -1.0,
                "tag": 0,
                "metaData": f"{audio_path.name};",
                "parameters": {},
            }],
            "metadata": {"timeAdded": _time_added_timestamp()},
        })
        source_id += 1

        image_source_ids.append(image_source_id)
        audio_source_ids.append(audio_source_id)

    # Add medias to timeline
    for i in range(pair_count):
        audio_info = audio_infos[i]
        audio_duration = float(audio_info["duration"])
        display_duration = max(float(min_image_duration), audio_duration)

        image_ticks = _seconds_to_ticks(display_duration)
        audio_ticks = _seconds_to_ticks(audio_duration)

        image_track_medias.append({
            "id": media_id,
            "_type": "IMFile",
            "src": image_source_ids[i],
            "trackNumber": 0,
            "trimStartSum": 0,
            "attributes": {"ident": f"slide_{i + 1}"},
            "parameters": {
                "geometryCrop0": 0.0,
                "geometryCrop1": 0.0,
                "geometryCrop2": 0.0,
                "geometryCrop3": 0.0,
            },
            "effects": [],
            "start": start_ticks,
            "duration": image_ticks,
            "mediaStart": 0,
            "mediaDuration": 1,
            "scalar": 1,
            "metadata": {
                "audiateLinkedSession": "",
                "clipSpeedAttribute": {"type": "bool", "value": False},
                "colorAttribute": {"type": "color", "value": [0, 0, 0, 0]},
                "default-height": {"type": "double", "value": float(project_height)},
                "default-translation0": {"type": "double", "value": 0.0},
                "default-translation1": {"type": "double", "value": 0.0},
                "default-translation2": {"type": "double", "value": 0.0},
                "default-width": {"type": "double", "value": float(project_width)},
            },
            "animationTracks": {},
        })
        media_id += 1

        audio_track_medias.append({
            "id": media_id,
            "_type": "AMFile",
            "src": audio_source_ids[i],
            "trackNumber": 0,
            "attributes": {
                "ident": f"audio_{i + 1}",
                "gain": 1.0,
                "mixToMono": False,
                "loudnessNormalization": True,
                "sourceFileOffset": 0,
            },
            "channelNumber": "0" if int(audio_info["num_channels"]) == 1 else "0,1",
            "parameters": {},
            "effects": [],
            "start": start_ticks,
            "duration": audio_ticks,
            "mediaStart": 0,
            "mediaDuration": audio_ticks,
            "scalar": 1,
            "metadata": {
                "audiateLinkedSession": "",
                "clipSpeedAttribute": {"type": "bool", "value": False},
                "colorAttribute": {"type": "color", "value": [0, 0, 0, 0]},
            },
            "animationTracks": {},
        })
        media_id += 1
        start_ticks += image_ticks
    
    # Write to file
    with open(output_path, 'w', encoding='utf-8-sig') as f:
        json.dump(project, f, indent=2)
    
    print(f"Created Camtasia project with audio: {output_path}")
    print(f"  - Image/audio pairs: {pair_count}")
    print(f"  - Resolution: {project_width}x{project_height}")
    print(f"  - Framerate: {int(round(framerate))} fps")
    if len(images) != len(audio_files):
        print(f"  - Note: {len(images)} images but {len(audio_files)} audio files (using {pair_count} pairs)")
    
    return output_path


# Keep old XML functions for backwards compatibility (deprecated)
def create_camtasia_project(
    image_folder: str,
    output_path: str = None,
    image_duration: float = 5.0,
    project_width: int = 1920,
    project_height: int = 1080,
    framerate: float = 30.0
) -> str:
    """
    DEPRECATED: Use create_camtasia_project_json() instead.
    Creates Camtasia 2025 project using new JSON format.
    """
    return create_camtasia_project_json(
        image_folder, output_path, image_duration, project_width, project_height, framerate
    )


def create_camtasia_project_with_audio(
    image_folder: str,
    audio_folder: str,
    output_path: str = None,
    min_image_duration: float = 2.0,
    project_width: int = 1920,
    project_height: int = 1080,
    framerate: float = 30.0
) -> str:
    """
    DEPRECATED: Use create_camtasia_project_with_audio_json() instead.
    Creates Camtasia 2025 project using new JSON format.
    """
    return create_camtasia_project_with_audio_json(
        image_folder, audio_folder, output_path, min_image_duration, project_width, project_height, framerate
    )


def generate_camtasia_import_script(
    image_folder: str,
    output_path: str = None
) -> str:
    """
    Generate a text file with instructions and file paths for manual Camtasia import.
    Also creates a batch file to copy paths to clipboard.
    
    Args:
        image_folder: Path to folder containing images
        output_path: Path for output files
        
    Returns:
        Path to created instruction file
    """
    image_folder = Path(image_folder)
    images = get_images_from_directory(image_folder)
    
    if not images:
        raise ValueError(f"No images found in {image_folder}")
    
    if output_path is None:
        output_path = str(image_folder / "camtasia_import.txt")
    
    # Create instruction file
    with open(output_path, 'w') as f:
        f.write("=" * 60 + "\n")
        f.write("CAMTASIA IMPORT INSTRUCTIONS\n")
        f.write("=" * 60 + "\n\n")
        
        f.write(f"Folder: {image_folder.absolute()}\n")
        f.write(f"Images found: {len(images)}\n\n")
        
        f.write("METHOD 1 - Drag and Drop:\n")
        f.write("-" * 40 + "\n")
        f.write("1. Open Camtasia\n")
        f.write("2. Open this folder in File Explorer\n")
        f.write("3. Select all images (Ctrl+A)\n")
        f.write("4. Drag them to the Camtasia Media Bin\n")
        f.write("5. Then drag from Media Bin to Timeline\n\n")
        
        f.write("METHOD 2 - File Import:\n")
        f.write("-" * 40 + "\n")
        f.write("1. In Camtasia: File > Import > Media\n")
        f.write("2. Navigate to the folder listed above\n")
        f.write("3. Select all images and click Open\n\n")
        
        f.write("METHOD 3 - Copy Paths:\n")
        f.write("-" * 40 + "\n")
        f.write("Run 'copy_paths_to_clipboard.bat' to copy all paths\n\n")
        
        f.write("IMAGE FILES:\n")
        f.write("-" * 40 + "\n")
        for img in images:
            f.write(f"{img.absolute()}\n")
    
    # Create batch file to copy paths to clipboard
    batch_path = image_folder / "copy_paths_to_clipboard.bat"
    with open(batch_path, 'w') as f:
        f.write("@echo off\n")
        f.write("(\n")
        for img in images:
            f.write(f'echo "{img.absolute()}"\n')
        f.write(") | clip\n")
        f.write(f"echo Copied {len(images)} file paths to clipboard!\n")
        f.write("pause\n")
    
    # Create PowerShell script for more reliable clipboard
    ps_path = image_folder / "copy_paths_to_clipboard.ps1"
    with open(ps_path, 'w') as f:
        f.write("# Copy image paths to clipboard\n")
        f.write("$paths = @(\n")
        for img in images:
            f.write(f'    "{img.absolute()}"\n')
        f.write(")\n")
        f.write("$paths -join \"`n\" | Set-Clipboard\n")
        f.write(f'Write-Host "Copied {len(images)} file paths to clipboard!"')
        f.write("\nRead-Host 'Press Enter to exit'\n")
    
    print(f"Created Camtasia import helper files:")
    print(f"  - Instructions: {output_path}")
    print(f"  - Batch file: {batch_path}")
    print(f"  - PowerShell: {ps_path}")
    
    return output_path


def interactive_camtasia_export():
    """Interactive mode for creating Camtasia 2025 projects."""
    print("\n" + "=" * 60)
    print("CAMTASIA 2025 PROJECT GENERATOR")
    print("=" * 60)
    
    # Get folder
    folder = input("\nEnter path to image folder: ").strip()
    if not folder:
        print("No folder specified.")
        return
    
    folder = Path(folder)
    if not folder.exists():
        print(f"Folder not found: {folder}")
        return
    
    try:
        images = get_images_from_directory(folder)
        print(f"\nFound {len(images)} images")
    except Exception as e:
        print(f"Error: {e}")
        return
    
    print("\nGeneration options:")
    print("  1. Create Camtasia project (images only)")
    print("  2. Create Camtasia project with audio (pairs images with audio files)")
    print("  3. Create Audiate narration script (from PowerPoint speaker notes)")
    print("  4. Generate import helper files (instructions + scripts)")
    print("  5. All of the above")
    
    choice = input("\nEnter choice (1-5): ").strip()
    
    if choice in ['1', '5']:
        duration = input("Duration per image in seconds (default: 5): ").strip()
        duration = float(duration) if duration else 5.0
        
        try:
            create_camtasia_project_json(folder, image_duration=duration)
        except Exception as e:
            print(f"Error creating project: {e}")
    
    if choice in ['2', '5']:
        audio_folder = input("Enter path to audio folder (for option 2): ").strip()
        if audio_folder:
            audio_folder = Path(audio_folder)
            if audio_folder.exists():
                try:
                    create_camtasia_project_with_audio_json(folder, str(audio_folder))
                except Exception as e:
                    print(f"Error creating audio project: {e}")
            else:
                print(f"Audio folder not found: {audio_folder}")
    
    if choice in ['3', '5']:
        pptx_file = input("Enter path to PowerPoint file (for option 3): ").strip()
        if pptx_file:
            pptx_file = Path(pptx_file)
            if pptx_file.exists():
                try:
                    create_audiate_script_from_pptx(str(pptx_file))
                except Exception as e:
                    print(f"Error creating Audiate script: {e}")
            else:
                print(f"PowerPoint file not found: {pptx_file}")
    
    if choice in ['4', '5']:
        try:
            generate_camtasia_import_script(folder)
        except Exception as e:
            print(f"Error creating import scripts: {e}")
    
    print("\nDone!")


if __name__ == "__main__":
    import argparse
    
    parser = argparse.ArgumentParser(description="Camtasia Project Generator")
    parser.add_argument("folder", nargs='?', help="Path to image folder")
    parser.add_argument("-o", "--output", help="Output path for project file")
    parser.add_argument("-d", "--duration", type=float, default=5.0, 
                        help="Duration per image in seconds (default: 5)")
    parser.add_argument("--width", type=int, default=1920, help="Project width")
    parser.add_argument("--height", type=int, default=1080, help="Project height")
    parser.add_argument("--scripts-only", action="store_true",
                        help="Only generate import helper scripts (no project file)")
    
    args = parser.parse_args()
    
    if args.folder:
        if args.scripts_only:
            generate_camtasia_import_script(args.folder, args.output)
        else:
            create_camtasia_project(
                args.folder,
                args.output,
                args.duration,
                args.width,
                args.height
            )
            generate_camtasia_import_script(args.folder)
    else:
        interactive_camtasia_export()
