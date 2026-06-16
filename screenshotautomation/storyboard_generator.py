"""
Storyboard Generator - Create PowerPoint storyboards from captured screenshots.

Follows the Progress Software video storyboard template format:
  - Slide 1: Title slide (from template, with video title text updated)
  - Slides 2..N: One full-page screenshot per slide, with narration placeholder in notes
Output is a .pptx file ready for narration scripting and conversion to video.
"""

import sys
from pathlib import Path
from typing import Optional, List, Callable

try:
    from pptx import Presentation
    from pptx.util import Inches, Pt, Emu
    from pptx.enum.text import PP_ALIGN
    from pptx.dml.color import RGBColor
except ImportError as e:
    print(f"Missing required package: {e}")
    print("Please install: pip install python-pptx")
    sys.exit(1)


# Default template path — the Progress Software storyboard template
DEFAULT_TEMPLATE = (
    r"C:\Users\plangley.PROGRESS\Progress Software Corporation"
    r"\SP-Engineering - Video Team\Templates\storyboard sample nobrand.pptx"
)

IMAGE_EXTENSIONS = {".png", ".jpg", ".jpeg", ".bmp", ".gif", ".tiff", ".webp"}


class StoryboardGenerator:
    """Generate a PowerPoint storyboard from a folder of screenshots."""

    def __init__(
        self,
        template_path: Optional[str] = None,
        progress_callback: Optional[Callable[[str, int, int], None]] = None,
    ):
        """
        Args:
            template_path: Path to the .pptx template. Falls back to DEFAULT_TEMPLATE,
                           then to a plain 16:9 presentation if neither exists.
            progress_callback: Optional callable(message, current, total) for progress.
        """
        self.template_path = template_path or DEFAULT_TEMPLATE
        self.progress_callback = progress_callback

    # ------------------------------------------------------------------
    # Internal helpers
    # ------------------------------------------------------------------

    def _report(self, msg: str, current: int = 0, total: int = 0) -> None:
        if self.progress_callback:
            self.progress_callback(msg, current, total)
        else:
            print(msg)

    def get_image_files(self, folder: str) -> List[Path]:
        """Return sorted list of image files in *folder*."""
        return sorted(
            f
            for f in Path(folder).iterdir()
            if f.is_file() and f.suffix.lower() in IMAGE_EXTENSIONS
        )

    def _remove_slide(self, prs: Presentation, index: int) -> None:
        """Remove the slide at *index* from *prs*, including its package part."""
        from pptx.oxml.ns import qn as _qn
        sldIdLst = prs.slides._sldIdLst
        sldId = sldIdLst[index]
        # Retrieve the relationship id and drop both the rel and the slide part
        r_id = sldId.get(_qn("r:id"))
        if r_id:
            try:
                prs.part.drop_rel(r_id)
            except Exception:
                pass  # best-effort; proceed to remove from list anyway
        sldIdLst.remove(sldId)

    def _find_blank_layout(self, prs: Presentation):
        """Return the 'Blank' slide layout, or the last layout as fallback."""
        for layout in prs.slide_layouts:
            if layout.name.lower() == "blank":
                return layout
        return prs.slide_layouts[-1]

    def _set_notes(self, slide, image_path: Path, narration: str) -> None:
        """Set speaker notes with file path first, then bold narration."""
        notes_slide = slide.notes_slide
        tf = notes_slide.notes_text_frame
        tf.clear()

        # First line: absolute image file path.
        p1 = tf.paragraphs[0]
        p1.text = str(image_path.resolve())
        p1.alignment = PP_ALIGN.LEFT
        for run in p1.runs:
            run.font.size = Pt(11)
            run.font.bold = False
            run.font.color.rgb = RGBColor(0x00, 0x00, 0x00)

        # Second line: narration in bold.
        p2 = tf.add_paragraph()
        p2.text = narration.strip() if narration.strip() else "[No narration provided]"
        p2.alignment = PP_ALIGN.LEFT
        for run in p2.runs:
            run.font.size = Pt(13)
            run.font.bold = True
            run.font.color.rgb = RGBColor(0x00, 0x00, 0x00)

    def _add_screenshot_slide(
        self,
        prs: Presentation,
        blank_layout,
        image_path: Path,
        step_num: int,
        narration: str = "",
    ):
        """Add one screenshot slide to *prs* and return it."""
        slide = prs.slides.add_slide(blank_layout)

        slide_w = prs.slide_width
        slide_h = prs.slide_height

        # Screenshot fills the entire slide
        slide.shapes.add_picture(
            str(image_path),
            left=Emu(0),
            top=Emu(0),
            width=slide_w,
            height=slide_h,
        )

        # Speaker notes: file path then bold narration.
        self._set_notes(
            slide,
            image_path=image_path,
            narration=narration or f"[Step {step_num} narration goes here]",
        )

        return slide

    def _build_narration(self, template: str, step_num: int, image_path: Path) -> str:
        """Build narration text for speaker notes from a format template."""
        try:
            resolved_path = image_path.resolve()
            return template.format(
                step=step_num,
                file_name=image_path.name,
                file_stem=image_path.stem,
                file_path=str(resolved_path),
                folder_path=str(resolved_path.parent),
            )
        except Exception:
            # Fallback if user-provided template has invalid placeholders
            return (
                f"Step {step_num}: [Add narration for this step]\n"
                f"Image path: {image_path.resolve()}"
            )

    # ------------------------------------------------------------------
    # Public API
    # ------------------------------------------------------------------

    def generate(
        self,
        screenshots_folder: str,
        output_path: str,
        video_title: str = "Video Title",
        narration_template: str = (
            "Step {step}: [Add narration for this step]\n"
            "Image path: {file_path}"
        ),
    ) -> int:
        """
        Generate a storyboard PPTX from a folder of screenshots.

        Args:
            screenshots_folder: Folder containing the screenshot images.
            output_path:        Destination .pptx file path.
            video_title:        Text to place on the title slide.
            narration_template: Format string for speaker-note narration.
                                Supported tokens: {step}, {file_name}, {file_stem},
                                {file_path}, {folder_path}

        Returns:
            The number of screenshot slides added.

        Raises:
            ValueError: If no images are found in *screenshots_folder*.
        """
        images = self.get_image_files(screenshots_folder)
        if not images:
            raise ValueError(f"No images found in: {screenshots_folder}")

        total_steps = len(images) + 2  # open/save + one per image

        # --- Open (or build) the base presentation ---
        self._report("Opening template...", 0, total_steps)

        if Path(self.template_path).exists():
            prs = Presentation(self.template_path)

            # Keep only slide 0 (title slide); discard all others so we can
            # append screenshot slides cleanly after it.
            slide_count = len(prs.slides)
            for i in range(slide_count - 1, 0, -1):
                self._remove_slide(prs, i)

            # Update the title text on the retained title slide.
            # The template has a TextBox whose text reads 'Title' — find it
            # and overwrite it with the user-supplied video title.
            title_slide = prs.slides[0]
            for shape in title_slide.shapes:
                if not shape.has_text_frame:
                    continue
                tf = shape.text_frame
                for para in tf.paragraphs:
                    if para.text.strip():
                        # Clear all runs, write new text into first run
                        for run in para.runs:
                            run.text = ""
                        if para.runs:
                            para.runs[0].text = video_title
                        else:
                            run = para.add_run()
                            run.text = video_title
                        break
                else:
                    continue
                break  # stop after the first shape with text

        else:
            # Fallback: build a minimal 16:9 presentation without branding
            prs = Presentation()
            prs.slide_width = Inches(13.333)
            prs.slide_height = Inches(7.5)

            blank_layout = self._find_blank_layout(prs)
            title_slide = prs.slides.add_slide(blank_layout)
            txBox = title_slide.shapes.add_textbox(
                Inches(0.5), Inches(3.0), Inches(12.333), Inches(1.5)
            )
            tf = txBox.text_frame
            para = tf.paragraphs[0]
            run = para.add_run()
            run.text = video_title
            run.font.size = Pt(40)
            run.font.bold = True
            run.font.color.rgb = RGBColor(0x1F, 0x49, 0x7D)

        # --- Add one slide per screenshot ---
        blank_layout = self._find_blank_layout(prs)

        for i, image_path in enumerate(images, 1):
            self._report(
                f"Adding slide {i} of {len(images)}: {image_path.name}",
                i,
                total_steps,
            )
            narration = self._build_narration(narration_template, i, image_path)
            self._add_screenshot_slide(prs, blank_layout, image_path, i, narration)

        # --- Save ---
        self._report("Saving storyboard...", total_steps - 1, total_steps)
        out = Path(output_path)
        out.parent.mkdir(parents=True, exist_ok=True)
        prs.save(str(out))
        self._report(f"Storyboard saved: {out}", total_steps, total_steps)

        return len(images)
