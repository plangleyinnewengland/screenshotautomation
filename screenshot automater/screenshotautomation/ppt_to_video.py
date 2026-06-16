"""
PowerPoint to Video Converter - Convert PPT/PPTX storyboards to narrated MP4 videos
Extracts slides as images and speaker notes as narration.
"""

import os
import sys
import tempfile
import shutil
import time
from pathlib import Path
from typing import Optional, List, Dict, Callable

try:
    from pptx import Presentation
    from pptx.util import Inches, Pt
except ImportError:
    print("Error: python-pptx not installed. Run: pip install python-pptx")
    sys.exit(1)

try:
    from PIL import Image
    from moviepy import (
        ImageClip, AudioFileClip, concatenate_videoclips, ColorClip
    )
    import pyttsx3
except ImportError as e:
    print(f"Missing required package: {e}")
    print("Please install: pip install pillow moviepy pyttsx3")
    sys.exit(1)

# For converting PPT slides to images, we need COM automation on Windows
import platform
if platform.system() == 'Windows':
    try:
        import comtypes.client
        HAS_COMTYPES = True
    except ImportError:
        HAS_COMTYPES = False
        print("Warning: comtypes not installed. Install with: pip install comtypes")
else:
    HAS_COMTYPES = False


class PPTToVideoConverter:
    """Convert PowerPoint presentations to narrated MP4 videos."""
    
    def __init__(
        self,
        output_resolution: tuple = (1920, 1080),
        fps: int = 30,
        default_slide_duration: float = 5.0,
        tts_rate: int = 150,
        tts_volume: float = 1.0,
        progress_callback: Optional[Callable[[str, int, int], None]] = None
    ):
        """
        Initialize the converter.
        
        Args:
            output_resolution: Video resolution (width, height)
            fps: Frames per second
            default_slide_duration: Default duration per slide if no narration
            tts_rate: Text-to-speech rate (words per minute)
            tts_volume: TTS volume (0.0 to 1.0)
            progress_callback: Optional callback(status, current, total) for progress updates
        """
        self.output_resolution = output_resolution
        self.fps = fps
        self.default_slide_duration = default_slide_duration
        self.tts_rate = tts_rate
        self.tts_volume = tts_volume
        self.progress_callback = progress_callback
        self._selected_voice_id = None
        
        # Initialize TTS
        self._init_tts()
        
    def _init_tts(self):
        """Initialize text-to-speech engine."""
        try:
            self.tts_engine = pyttsx3.init()
            self.tts_engine.setProperty('rate', self.tts_rate)
            self.tts_engine.setProperty('volume', self.tts_volume)
            
            voices = self.tts_engine.getProperty('voices')
            self.available_voices = [(v.id, v.name) for v in voices]
        except Exception as e:
            print(f"Warning: Could not initialize TTS: {e}")
            self.tts_engine = None
            self.available_voices = []
            
    def list_voices(self) -> List[tuple]:
        """List available TTS voices as (id, name) tuples."""
        return self.available_voices
    
    def set_voice(self, voice_id: str):
        """Set the TTS voice by ID."""
        self._selected_voice_id = voice_id
        if self.tts_engine:
            self.tts_engine.setProperty('voice', voice_id)
            
    def set_tts_rate(self, rate: int):
        """Set TTS speech rate."""
        self.tts_rate = rate
        if self.tts_engine:
            self.tts_engine.setProperty('rate', rate)
            
    def _report_progress(self, status: str, current: int, total: int):
        """Report progress if callback is set."""
        if self.progress_callback:
            self.progress_callback(status, current, total)
        else:
            print(f"{status} ({current}/{total})")
            
    def extract_slides_info(self, pptx_path: str) -> List[Dict]:
        """
        Extract slide information from a PowerPoint file.
        
        Args:
            pptx_path: Path to the .pptx file
            
        Returns:
            List of dicts with 'slide_number' and 'notes' for each slide
        """
        prs = Presentation(pptx_path)
        slides_info = []
        
        for i, slide in enumerate(prs.slides):
            notes_text = ""
            if slide.has_notes_slide:
                notes_slide = slide.notes_slide
                notes_text = notes_slide.notes_text_frame.text.strip()
            
            slides_info.append({
                'slide_number': i + 1,
                'notes': notes_text
            })
            
        return slides_info
    
    def export_slides_as_images(self, pptx_path: str, output_dir: str) -> List[Path]:
        """
        Export PowerPoint slides as images.
        
        Args:
            pptx_path: Path to the .pptx file
            output_dir: Directory to save images
            
        Returns:
            List of paths to exported images
        """
        pptx_path = Path(pptx_path).resolve()
        output_dir = Path(output_dir)
        output_dir.mkdir(parents=True, exist_ok=True)
        
        if platform.system() == 'Windows' and HAS_COMTYPES:
            return self._export_slides_com(pptx_path, output_dir)
        else:
            # Fallback: try using LibreOffice if available
            return self._export_slides_libreoffice(pptx_path, output_dir)
    
    def _export_slides_com(self, pptx_path: Path, output_dir: Path) -> List[Path]:
        """Export slides using PowerPoint COM automation (Windows only)."""
        image_paths = []
        powerpoint = None
        presentation = None
        
        try:
            # Initialize PowerPoint
            powerpoint = comtypes.client.CreateObject("PowerPoint.Application")
            powerpoint.Visible = 1
            
            # Open presentation
            presentation = powerpoint.Presentations.Open(str(pptx_path), WithWindow=False)
            
            total_slides = presentation.Slides.Count
            
            for i in range(1, total_slides + 1):
                self._report_progress(f"Exporting slide {i}", i, total_slides)
                
                # Export slide as PNG
                output_path = output_dir / f"slide_{i:03d}.png"
                presentation.Slides(i).Export(str(output_path), "PNG", 
                                              self.output_resolution[0], 
                                              self.output_resolution[1])
                
                if output_path.exists():
                    image_paths.append(output_path)
                    
        except Exception as e:
            print(f"Error exporting slides via COM: {e}")
            raise
        finally:
            if presentation:
                presentation.Close()
            if powerpoint:
                powerpoint.Quit()
                
        return image_paths
    
    def _export_slides_libreoffice(self, pptx_path: Path, output_dir: Path) -> List[Path]:
        """Export slides using LibreOffice (cross-platform fallback)."""
        import subprocess
        
        # Try to find LibreOffice
        soffice_paths = [
            r"C:\Program Files\LibreOffice\program\soffice.exe",
            r"C:\Program Files (x86)\LibreOffice\program\soffice.exe",
            "/usr/bin/soffice",
            "/usr/bin/libreoffice",
        ]
        
        soffice = None
        for path in soffice_paths:
            if Path(path).exists():
                soffice = path
                break
                
        if not soffice:
            raise RuntimeError(
                "Could not find PowerPoint or LibreOffice to export slides. "
                "Please install Microsoft PowerPoint or LibreOffice."
            )
        
        # Export to PDF first, then convert to images
        self._report_progress("Converting with LibreOffice...", 0, 1)
        
        # Export as images directly
        subprocess.run([
            soffice,
            "--headless",
            "--convert-to", "png",
            "--outdir", str(output_dir),
            str(pptx_path)
        ], check=True)
        
        # Find exported images
        image_paths = sorted(output_dir.glob("*.png"))
        
        if not image_paths:
            raise RuntimeError("No images were exported from the presentation")
            
        return image_paths
    
    def _generate_audio(self, text: str, output_path: str) -> Optional[str]:
        """Generate audio file from text using TTS with timeout protection."""
        if not text.strip():
            return None
        
        import multiprocessing
        
        def _tts_worker(text, output_path, rate, volume, voice_id):
            """Worker function to run TTS in separate process."""
            try:
                import pyttsx3
                engine = pyttsx3.init()
                engine.setProperty('rate', rate)
                engine.setProperty('volume', volume)
                if voice_id:
                    engine.setProperty('voice', voice_id)
                engine.save_to_file(text, output_path)
                engine.runAndWait()
                engine.stop()
            except Exception as e:
                print(f"TTS worker error: {e}")
        
        try:
            # Run TTS in a separate process with timeout to prevent hanging
            process = multiprocessing.Process(
                target=_tts_worker,
                args=(text, output_path, self.tts_rate, self.tts_volume, self._selected_voice_id)
            )
            process.start()
            
            # Wait up to 60 seconds for audio generation
            process.join(timeout=60)
            
            if process.is_alive():
                print(f"Warning: TTS timed out, terminating...")
                process.terminate()
                process.join(timeout=5)
                return None
            
            time.sleep(0.2)
            
            if Path(output_path).exists() and Path(output_path).stat().st_size > 0:
                return output_path
                
        except Exception as e:
            print(f"Warning: Could not generate audio: {e}")
            
        return None
    
    def _resize_image(self, image_path: Path, target_size: tuple) -> Image.Image:
        """Resize image to target size with letterboxing."""
        img = Image.open(image_path)
        img_ratio = img.width / img.height
        target_ratio = target_size[0] / target_size[1]
        
        if img_ratio > target_ratio:
            new_width = target_size[0]
            new_height = int(new_width / img_ratio)
        else:
            new_height = target_size[1]
            new_width = int(new_height * img_ratio)
            
        img_resized = img.resize((new_width, new_height), Image.Resampling.LANCZOS)
        
        background = Image.new('RGB', target_size, (0, 0, 0))
        paste_x = (target_size[0] - new_width) // 2
        paste_y = (target_size[1] - new_height) // 2
        
        if img_resized.mode == 'RGBA':
            background.paste(img_resized, (paste_x, paste_y), img_resized)
        else:
            background.paste(img_resized, (paste_x, paste_y))
            
        return background
    
    def convert(
        self,
        pptx_path: str,
        output_path: str,
        custom_narrations: Optional[Dict[int, str]] = None
    ) -> str:
        """
        Convert a PowerPoint presentation to an MP4 video.
        
        Args:
            pptx_path: Path to the .pptx file
            output_path: Path for the output .mp4 file
            custom_narrations: Optional dict mapping slide numbers to custom narration text
            
        Returns:
            Path to the created video file
        """
        pptx_path = Path(pptx_path)
        output_path = Path(output_path)
        
        if not pptx_path.exists():
            raise FileNotFoundError(f"PowerPoint file not found: {pptx_path}")
            
        # Create temp directory for working files
        temp_dir = Path(tempfile.mkdtemp(prefix="ppt_to_video_"))
        slides_dir = temp_dir / "slides"
        audio_dir = temp_dir / "audio"
        slides_dir.mkdir()
        audio_dir.mkdir()
        
        try:
            # Step 1: Get slide notes
            self._report_progress("Reading presentation...", 0, 1)
            slides_info = self.extract_slides_info(pptx_path)
            total_slides = len(slides_info)
            
            # Step 2: Export slides as images
            self._report_progress("Exporting slides...", 0, total_slides)
            image_paths = self.export_slides_as_images(pptx_path, slides_dir)
            
            if len(image_paths) != total_slides:
                print(f"Warning: Expected {total_slides} slides but exported {len(image_paths)} images")
                
            # Step 3: Generate audio for each slide
            audio_paths = []
            for i, info in enumerate(slides_info):
                slide_num = info['slide_number']
                
                # Use custom narration if provided, otherwise use notes
                narration = info['notes']
                if custom_narrations and slide_num in custom_narrations:
                    narration = custom_narrations[slide_num]
                    
                self._report_progress(f"Generating audio for slide {slide_num}", i + 1, total_slides)
                
                audio_path = None
                if narration.strip():
                    audio_file = audio_dir / f"audio_{slide_num:03d}.wav"
                    audio_path = self._generate_audio(narration, str(audio_file))
                    
                audio_paths.append(audio_path)
                
            # Step 4: Create video clips for each slide
            clips = []
            temp_images = []
            
            for i, (image_path, audio_path) in enumerate(zip(image_paths, audio_paths)):
                slide_num = i + 1
                self._report_progress(f"Creating clip for slide {slide_num}", slide_num, total_slides)
                
                # Resize image
                resized = self._resize_image(image_path, self.output_resolution)
                temp_img = temp_dir / f"resized_{slide_num:03d}.png"
                resized.save(temp_img)
                temp_images.append(temp_img)
                
                # Determine duration
                duration = self.default_slide_duration
                audio_clip = None
                
                if audio_path and Path(audio_path).exists():
                    try:
                        audio_clip = AudioFileClip(audio_path)
                        # Add padding after narration
                        duration = max(duration, audio_clip.duration + 1.0)
                    except Exception as e:
                        print(f"Warning: Could not load audio for slide {slide_num}: {e}")
                        
                # Create image clip
                clip = ImageClip(str(temp_img), duration=duration)
                
                if audio_clip:
                    clip = clip.with_audio(audio_clip)
                    
                clips.append(clip)
                
            # Step 5: Concatenate all clips
            self._report_progress("Combining clips...", total_slides, total_slides)
            
            if not clips:
                raise ValueError("No clips were created")
                
            final_clip = concatenate_videoclips(clips, method="compose")
            
            # Step 6: Export video
            self._report_progress("Exporting video...", 0, 1)
            output_path.parent.mkdir(parents=True, exist_ok=True)
            
            final_clip.write_videofile(
                str(output_path),
                fps=self.fps,
                codec='libx264',
                audio_codec='aac',
                logger=None
            )
            
            # Clean up clips
            final_clip.close()
            for clip in clips:
                clip.close()
                
            self._report_progress("Complete!", 1, 1)
            return str(output_path)
            
        finally:
            # Clean up temp directory
            try:
                shutil.rmtree(temp_dir)
            except Exception as e:
                print(f"Warning: Could not clean up temp files: {e}")
                
    def preview_notes(self, pptx_path: str) -> List[Dict]:
        """
        Preview the slides and notes from a PowerPoint file.
        
        Args:
            pptx_path: Path to the .pptx file
            
        Returns:
            List of dicts with slide info including notes
        """
        return self.extract_slides_info(pptx_path)


def main():
    """Command-line interface."""
    import argparse
    
    parser = argparse.ArgumentParser(
        description="Convert PowerPoint presentations to narrated MP4 videos"
    )
    parser.add_argument("input", help="Path to the PowerPoint file (.pptx)")
    parser.add_argument("-o", "--output", help="Output video path (default: input name with .mp4)")
    parser.add_argument("--resolution", default="1920x1080", help="Video resolution (default: 1920x1080)")
    parser.add_argument("--fps", type=int, default=30, help="Frames per second (default: 30)")
    parser.add_argument("--duration", type=float, default=5.0, help="Default slide duration in seconds")
    parser.add_argument("--tts-rate", type=int, default=150, help="TTS speech rate (default: 150)")
    parser.add_argument("--preview", action="store_true", help="Preview notes without creating video")
    parser.add_argument("--list-voices", action="store_true", help="List available TTS voices")
    parser.add_argument("--voice", help="TTS voice ID to use")
    
    args = parser.parse_args()
    
    # Parse resolution
    try:
        width, height = map(int, args.resolution.lower().split('x'))
        resolution = (width, height)
    except ValueError:
        print(f"Invalid resolution format: {args.resolution}")
        sys.exit(1)
        
    converter = PPTToVideoConverter(
        output_resolution=resolution,
        fps=args.fps,
        default_slide_duration=args.duration,
        tts_rate=args.tts_rate
    )
    
    if args.list_voices:
        print("Available TTS voices:")
        for voice_id, voice_name in converter.list_voices():
            print(f"  {voice_name}")
            print(f"    ID: {voice_id}")
        sys.exit(0)
        
    if args.voice:
        converter.set_voice(args.voice)
        
    if args.preview:
        print(f"Preview of: {args.input}\n")
        slides = converter.preview_notes(args.input)
        for slide in slides:
            print(f"Slide {slide['slide_number']}:")
            if slide['notes']:
                print(f"  Notes: {slide['notes'][:100]}{'...' if len(slide['notes']) > 100 else ''}")
            else:
                print("  (No notes)")
            print()
        sys.exit(0)
        
    # Determine output path
    output_path = args.output
    if not output_path:
        input_path = Path(args.input)
        output_path = input_path.with_suffix('.mp4')
        
    print(f"Converting: {args.input}")
    print(f"Output: {output_path}")
    print()
    
    result = converter.convert(args.input, output_path)
    print(f"\nVideo created: {result}")


if __name__ == "__main__":
    main()
