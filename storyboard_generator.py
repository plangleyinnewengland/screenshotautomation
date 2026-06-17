"""
Storyboard Generator - Create PowerPoint storyboards from captured screenshots.

Follows the Progress Software video storyboard template format:
  - Slide 1: Title slide (from template, with video title text updated)
  - Slides 2..N: One full-page screenshot per slide, with narration placeholder in notes
Output is a .pptx file ready for narration scripting and conversion to video.
"""

import sys
from pathlib import Path
from typing import Optional, List, Callable

try:
    from pptx import Presentation
    from pptx.util import Inches, Pt, Emu
    from pptx.enum.text import PP_ALIGN
    from pptx.dml.color import RGBColor
except ImportError as e:
    print(f"Missing required package: {e}")
    print("Please install: pip install python-pptx")
    sys.exit(1)


# Default template path — the Progress Software storyboard template
DEFAULT_TEMPLATE = (
    r"C:\Users\plangley.PROGRESS\Progress Software Corporation"
    r"\SP-Engineering - Video Team\Templates\storyboard sample nobrand.pptx"
)

# Default title page image (slide 1)
DEFAULT_TITLE_PAGE = (
    r"C:\Users\plangley.PROGRESS\Progress Software Corporation"
    r"\SP-Engineering - Video Team\Templates\Slide1.PNG"
)

# Default ending slide image (more information slide)
DEFAULT_ENDING_SLIDE = (
    r"C:\Users\plangley.PROGRESS\Progress Software Corporation"
    r"\SP-Engineering - Video Team\Templates\Slide4.PNG"
)

# Default intro music for the title slide
DEFAULT_INTRO_MUSIC = (
    r"C:\Users\plangley.PROGRESS\Progress Software Corporation"
    r"\SP-Engineering - Video Team\Templates\intro_music.wav"
)

# Default narration for ending slide
DEFAULT_ENDING_NARRATION = "For additional resources, see the links in the Description."

IMAGE_EXTENSIONS = {".png", ".jpg", ".jpeg", ".bmp", ".gif", ".tiff", ".webp"}
VIDEO_EXTENSIONS = {".mp4"}
AUDIO_EXTENSIONS = {".wav", ".mp3"}


class StoryboardGenerator:
    """Generate a PowerPoint storyboard from a folder of screenshots."""

    def __init__(
        self,
        template_path: Optional[str] = None,
        progress_callback: Optional[Callable[[str, int, int], None]] = None,
    ):
        """
        Args:
            template_path: Path to the .pptx template. Falls back to DEFAULT_TEMPLATE,
                           then to a plain 16:9 presentation if neither exists.
            progress_callback: Optional callable(message, current, total) for progress.
        """
        self.template_path = template_path or DEFAULT_TEMPLATE
        self.progress_callback = progress_callback

    # ------------------------------------------------------------------
    # Internal helpers
    # ------------------------------------------------------------------

    def _report(self, msg: str, current: int = 0, total: int = 0) -> None:
        if self.progress_callback:
            self.progress_callback(msg, current, total)
        else:
            print(msg)

    def get_image_files(self, folder: str) -> List[Path]:
        """Return sorted list of image files in *folder*."""
        return sorted(
            f
            for f in Path(folder).iterdir()
            if f.is_file() and f.suffix.lower() in IMAGE_EXTENSIONS
        )

    def _remove_slide(self, prs: Presentation, index: int) -> None:
        """Remove the slide at *index* from *prs*, including its package part."""
        from pptx.oxml.ns import qn as _qn
        sldIdLst = prs.slides._sldIdLst
        sldId = sldIdLst[index]
        # Retrieve the relationship id and drop both the rel and the slide part
        r_id = sldId.get(_qn("r:id"))
        if r_id:
            try:
                prs.part.drop_rel(r_id)
            except Exception:
                pass  # best-effort; proceed to remove from list anyway
        sldIdLst.remove(sldId)

    def _find_blank_layout(self, prs: Presentation):
        """Return the 'Blank' slide layout, or the last layout as fallback."""
        for layout in prs.slide_layouts:
            if layout.name.lower() == "blank":
                return layout
        return prs.slide_layouts[-1]

    def _set_notes(self, slide, image_path: Path, narration: str) -> None:
        """Set speaker notes with file path first, then bold narration."""
        notes_slide = slide.notes_slide
        tf = notes_slide.notes_text_frame
        tf.clear()

        # First line: absolute image file path.
        p1 = tf.paragraphs[0]
        p1.text = str(image_path.resolve())
        p1.alignment = PP_ALIGN.LEFT
        for run in p1.runs:
            run.font.size = Pt(11)
            run.font.bold = False
            run.font.color.rgb = RGBColor(0x00, 0x00, 0x00)

        # Second line: narration in bold.
        p2 = tf.add_paragraph()
        p2.text = narration.strip() if narration.strip() else "[No narration provided]"
        p2.alignment = PP_ALIGN.LEFT
        for run in p2.runs:
            run.font.size = Pt(13)
            run.font.bold = True
            run.font.color.rgb = RGBColor(0x00, 0x00, 0x00)

    def _add_screenshot_slide(
        self,
        prs: Presentation,
        blank_layout,
        image_path: Path,
        step_num: int,
        narration: str = "",
    ):
        """Add one screenshot slide to *prs* and return it."""
        slide = prs.slides.add_slide(blank_layout)

        slide_w = prs.slide_width
        slide_h = prs.slide_height

        # Screenshot fills the entire slide
        slide.shapes.add_picture(
            str(image_path),
            left=Emu(0),
            top=Emu(0),
            width=slide_w,
            height=slide_h,
        )

        # Speaker notes: file path then bold narration.
        self._set_notes(
            slide,
            image_path=image_path,
            narration=narration or f"[Step {step_num} narration goes here]",
        )

        return slide

    def _add_ending_media_slide(
        self,
        prs: Presentation,
        blank_layout,
        media_path: Path,
        narration: str,
    ) -> None:
        """Add an ending slide using either an image or an embedded video."""
        slide = prs.slides.add_slide(blank_layout)
        slide_w = prs.slide_width
        slide_h = prs.slide_height
        suffix = media_path.suffix.lower()

        if suffix in VIDEO_EXTENSIONS:
            # Embed video directly so it can be played in PowerPoint.
            slide.shapes.add_movie(
                str(media_path),
                left=Emu(0),
                top=Emu(0),
                width=slide_w,
                height=slide_h,
                mime_type="video/mp4",
            )
        else:
            slide.shapes.add_picture(
                str(media_path),
                left=Emu(0),
                top=Emu(0),
                width=slide_w,
                height=slide_h,
            )

        self._set_notes(slide, media_path, narration)

    def _add_intro_audio(self, slide, audio_path: Path, slide_w: int, slide_h: int) -> None:
        """Embed intro audio on title slide as a small media object."""
        suffix = audio_path.suffix.lower()
        if suffix not in AUDIO_EXTENSIONS:
            return

        mime_type = "audio/wav" if suffix == ".wav" else "audio/mpeg"
        icon_w = Inches(0.25)
        icon_h = Inches(0.25)
        left = max(Emu(0), slide_w - icon_w - Inches(0.1))
        top = max(Emu(0), slide_h - icon_h - Inches(0.1))

        slide.shapes.add_movie(
            str(audio_path),
            left=left,
            top=top,
            width=icon_w,
            height=icon_h,
            mime_type=mime_type,
        )

    def _build_narration(self, template: str, step_num: int, image_path: Path) -> str:
        """Build narration text for speaker notes from a format template."""
        try:
            resolved_path = image_path.resolve()
            return template.format(
                step=step_num,
                file_name=image_path.name,
                file_stem=image_path.stem,
                file_path=str(resolved_path),
                folder_path=str(resolved_path.parent),
            )
        except Exception:
            # Fallback if user-provided template has invalid placeholders
            return (
                f"Step {step_num}: [Add narration for this step]\n"
                f"Image path: {image_path.resolve()}"
            )

    # ------------------------------------------------------------------
    # Public API
    # ------------------------------------------------------------------

    def generate(
        self,
        screenshots_folder: str,
        output_path: str,
        video_title: str = "Video Title",
        narration_template: str = (
            "Step {step}: [Add narration for this step]\n"
            "Image path: {file_path}"
        ),
        title_page_image: Optional[str] = None,
        intro_music_path: Optional[str] = None,
        ending_slide_image: Optional[str] = None,
        ending_slide_narration: Optional[str] = None,
    ) -> int:
        """
        Generate a storyboard PPTX from a folder of screenshots.

        Args:
            screenshots_folder: Folder containing the screenshot images.
            output_path:        Destination .pptx file path.
            video_title:        Text to place on the title slide.
            narration_template: Format string for speaker-note narration.
                                Supported tokens: {step}, {file_name}, {file_stem},
                                {file_path}, {folder_path}
            title_page_image:   Path to custom title page image. If None, uses DEFAULT_TITLE_PAGE.
                                The video title text will be overlaid on this image.
            intro_music_path:   Path to intro music for title slide. Supports .wav and .mp3.
                                If None, uses DEFAULT_INTRO_MUSIC.
            ending_slide_image: Path to custom ending slide media. Supports images and .mp4.
                                If None, uses DEFAULT_ENDING_SLIDE.
            ending_slide_narration: Narration for the ending slide. If None, uses DEFAULT_ENDING_NARRATION.

        Returns:
            The number of screenshot slides added (excluding title and ending slides).

        Raises:
            ValueError: If no images are found in *screenshots_folder*.
        """
        images = self.get_image_files(screenshots_folder)
        if not images:
            raise ValueError(f"No images found in: {screenshots_folder}")

        # Determine which title page and ending slide to use
        title_page = title_page_image or DEFAULT_TITLE_PAGE
        intro_music = intro_music_path or DEFAULT_INTRO_MUSIC
        ending_slide = ending_slide_image or DEFAULT_ENDING_SLIDE
        ending_narration = ending_slide_narration or DEFAULT_ENDING_NARRATION

        total_steps = len(images) + 3  # title + screenshots + ending + save

        # --- Create a new presentation ---
        self._report("Creating presentation...", 0, total_steps)
        prs = Presentation()
        prs.slide_width = Inches(13.333)
        prs.slide_height = Inches(7.5)
        blank_layout = self._find_blank_layout(prs)

        # --- Add title page ---
        self._report("Adding title page...", 1, total_steps)
        title_slide = prs.slides.add_slide(blank_layout)
        slide_w = prs.slide_width
        slide_h = prs.slide_height

        # Add the title page image as background
        if Path(title_page).exists():
            title_slide.shapes.add_picture(
                str(title_page),
                left=Emu(0),
                top=Emu(0),
                width=slide_w,
                height=slide_h,
            )
            
            # Overlay the video title text in the center
            left = Inches(0.5)
            top = Inches(2.5)
            width = Inches(12.333)
            height = Inches(2.5)
            
            txBox = title_slide.shapes.add_textbox(left, top, width, height)
            tf = txBox.text_frame
            tf.word_wrap = True
            para = tf.paragraphs[0]
            para.alignment = PP_ALIGN.CENTER
            run = para.add_run()
            run.text = video_title
            run.font.size = Pt(54)
            run.font.bold = True
            run.font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)  # White text
        else:
            # Fallback: create a simple title slide
            txBox = title_slide.shapes.add_textbox(
                Inches(0.5), Inches(3.0), Inches(12.333), Inches(1.5)
            )
            tf = txBox.text_frame
            para = tf.paragraphs[0]
            para.alignment = PP_ALIGN.CENTER
            run = para.add_run()
            run.text = video_title
            run.font.size = Pt(40)
            run.font.bold = True
            run.font.color.rgb = RGBColor(0x1F, 0x49, 0x7D)

        if Path(intro_music).exists():
            self._add_intro_audio(title_slide, Path(intro_music), slide_w, slide_h)

        # --- Add one slide per screenshot ---
        for i, image_path in enumerate(images, 1):
            self._report(
                f"Adding slide {i} of {len(images)}: {image_path.name}",
                i + 1,
                total_steps,
            )
            narration = self._build_narration(narration_template, i, image_path)
            self._add_screenshot_slide(prs, blank_layout, image_path, i, narration)

        # --- Add ending slide ---
        self._report("Adding ending slide...", len(images) + 2, total_steps)
        if Path(ending_slide).exists():
            self._add_ending_media_slide(
                prs=prs,
                blank_layout=blank_layout,
                media_path=Path(ending_slide),
                narration=ending_narration,
            )
        # If ending slide media doesn't exist, silently skip it (no error)

        # --- Save ---
        self._report("Saving storyboard...", total_steps - 1, total_steps)
        out = Path(output_path)
        out.parent.mkdir(parents=True, exist_ok=True)
        prs.save(str(out))
        self._report(f"Storyboard saved: {out}", total_steps, total_steps)

        return len(images)
